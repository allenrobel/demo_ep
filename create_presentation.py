#!/usr/bin/env python3
"""Generate PowerPoint presentation for endpoint library design."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(24)
        sub_para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        sub_para.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add colored background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.2), Inches(10), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    title_para.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, code_example=None):
    """Add a content slide with bullets and optional code."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(9), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x00, 0x7a, 0xcc)
    line.line.fill.background()

    # Calculate content area based on whether we have code
    if code_example:
        bullet_height = 2.0
        bullet_top = 1.4
    else:
        bullet_height = 4.5
        bullet_top = 1.4

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(bullet_top), Inches(9), Inches(bullet_height))
    bullet_frame = bullet_box.text_frame
    bullet_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            para = bullet_frame.paragraphs[0]
        else:
            para = bullet_frame.add_paragraph()
        para.text = f"• {bullet}"
        para.font.size = Pt(20)
        para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        para.space_after = Pt(12)

    # Code example
    if code_example:
        code_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(2.5))
        code_frame = code_box.text_frame
        code_frame.word_wrap = True

        # Code background
        code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(3.5), Inches(9.2), Inches(2.6))
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)
        code_bg.line.color.rgb = RGBColor(0xdd, 0xdd, 0xdd)

        # Move code box to front
        spTree = slide.shapes._spTree
        code_bg_sp = code_bg._element
        code_box_sp = code_box._element
        spTree.remove(code_bg_sp)
        spTree.insert(2, code_bg_sp)

        code_para = code_frame.paragraphs[0]
        code_para.text = code_example
        code_para.font.size = Pt(14)
        code_para.font.name = "Courier New"
        code_para.font.color.rgb = RGBColor(0x2d, 0x2d, 0x2d)

    return slide

def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column comparison slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

    # Left column header
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(0.5))
    lh_frame = left_header.text_frame
    lh_para = lh_frame.paragraphs[0]
    lh_para.text = left_title
    lh_para.font.size = Pt(22)
    lh_para.font.bold = True
    lh_para.font.color.rgb = RGBColor(0xcc, 0x00, 0x00)

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.2), Inches(4))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            para = left_frame.paragraphs[0]
        else:
            para = left_frame.add_paragraph()
        para.text = f"• {item}"
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        para.space_after = Pt(8)

    # Right column header
    right_header = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.2), Inches(0.5))
    rh_frame = right_header.text_frame
    rh_para = rh_frame.paragraphs[0]
    rh_para.text = right_title
    rh_para.font.size = Pt(22)
    rh_para.font.bold = True
    rh_para.font.color.rgb = RGBColor(0x00, 0x7a, 0x00)

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.9), Inches(4.2), Inches(4))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            para = right_frame.paragraphs[0]
        else:
            para = right_frame.add_paragraph()
        para.text = f"• {item}"
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        para.space_after = Pt(8)

    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Typed Endpoint Library",
        "Design Goals & Benefits"
    )

    # Slide 2: The Question
    add_content_slide(
        prs,
        "The Question We're Addressing",
        [
            '"Why all this complexity?"',
            '"Why not just hardcode endpoint paths?"',
            '"Why not use a generic query-string builder?"',
            "Let's explore why this design pays dividends..."
        ]
    )

    # Slide 3: Section - Design Goals
    add_section_slide(prs, "Design Goals")

    # Slide 4: Type Safety
    add_content_slide(
        prs,
        "1. Type Safety as a First-Class Concern",
        [
            "Pydantic validates parameters on assignment",
            "Field constraints enforce valid ranges (min_length, max_length)",
            "Enums restrict values to API-valid options",
            "Errors caught at development time, not runtime"
        ],
        "fabric_name: Optional[str] = Field(\n    default=None,\n    min_length=1,\n    max_length=64\n)"
    )

    # Slide 5: Single Source of Truth
    add_content_slide(
        prs,
        "2. Single Source of Truth",
        [
            "Base URL paths → BasePath class",
            "HTTP verbs → Each endpoint's verb property",
            "Required parameters → Mixin fields + path validation",
            "Query parameter formats → to_query_string() methods",
            "When API changes, updates are localized"
        ]
    )

    # Slide 6: Composition Over Inheritance
    add_content_slide(
        prs,
        "3. Composition Over Inheritance",
        [
            "Mixins provide reusable field definitions",
            "Query parameters are composable",
            "No deep inheritance hierarchies",
            "Mix and match capabilities as needed"
        ],
        "class EpFabricDelete(FabricNameMixin, BaseModel):\n    # Inherits fabric_name with all validation\n    ..."
    )

    # Slide 7: Separation of Concerns
    add_content_slide(
        prs,
        "4. Separation of Concerns",
        [
            "base_paths.py → URL path construction",
            "query_params.py → Query string building & validation",
            "enums.py → Type-safe enumerated values",
            "onemanage_*_endpoints.py → Domain-specific endpoint modules",
            "Each module independently testable"
        ]
    )

    # Slide 8: Consistent Interface
    add_content_slide(
        prs,
        "5. Consistent Developer Interface",
        [
            "Every endpoint works the same way",
            "Learn one pattern, apply everywhere",
            "IDE autocomplete guides usage"
        ],
        "request = EpFabricConfigDeploy()\nrequest.fabric_name = \"MyFabric\"\nrequest.query_params.force_show_run = \"true\"\n\npath = request.path  # Complete URL\nverb = request.verb  # HTTP method"
    )

    # Slide 9: Section - Benefits
    add_section_slide(prs, "Key Benefits")

    # Slide 10: Comparison
    add_comparison_slide(
        prs,
        "Hardcoded Strings vs. Typed Endpoints",
        "Hardcoded Approach",
        [
            "Errors surface at runtime (404s, 400s)",
            "No IDE autocomplete",
            "Validation scattered or missing",
            "API changes require codebase search",
            "Easy to typo parameter names",
            "No documentation at point of use"
        ],
        "Typed Endpoint Approach",
        [
            "Errors caught before runtime",
            "Full IDE support & autocomplete",
            "Centralized validation rules",
            "API changes localized to domain modules",
            "Enums prevent typos",
            "Self-documenting with examples"
        ]
    )

    # Slide 11: Testability
    add_content_slide(
        prs,
        "Testability",
        [
            "Each endpoint class is independently instantiable",
            "Unit test path generation without HTTP calls",
            "Unit test parameter validation in isolation",
            "Unit test query string building separately",
            "Mocking is straightforward—pass objects, not strings"
        ]
    )

    # Slide 12: IDE Experience
    add_content_slide(
        prs,
        "Developer Experience",
        [
            "Autocomplete shows all available endpoints",
            "Hover for docstrings with usage examples",
            "Type checkers catch misuse immediately",
            "Refactoring tools rename safely across codebase",
            "New team members onboard faster"
        ]
    )

    # Slide 13: Maintainability
    add_content_slide(
        prs,
        "Maintainability at Scale",
        [
            "API contract encoded in the type system",
            "Change a constraint once, propagates everywhere",
            "Clear separation makes debugging easier",
            "Consistent patterns reduce cognitive load",
            "Documentation lives with the code"
        ]
    )

    # Slide 14: Graceful Degradation
    add_content_slide(
        prs,
        "Graceful Degradation",
        [
            "Works with or without Pydantic installed",
            "Fallback classes maintain basic functionality",
            "Supports diverse deployment environments",
            "Full validation when Pydantic available"
        ]
    )

    # Slide 15: Summary
    add_content_slide(
        prs,
        "Summary",
        [
            "Yes, it's more code upfront",
            "But: fewer runtime errors, better IDE support",
            "Centralized API knowledge = easier maintenance",
            "Type safety catches bugs before they ship",
            "Investment pays off as codebase scales"
        ]
    )

    # Slide 16: Closing
    add_title_slide(
        prs,
        "Questions?",
        "The complexity serves correctness, maintainability, and developer experience"
    )

    # Save
    prs.save("/Users/arobel/repos/demo_ep/endpoint_library_presentation.pptx")
    print("Presentation saved: endpoint_library_presentation.pptx")

if __name__ == "__main__":
    main()
